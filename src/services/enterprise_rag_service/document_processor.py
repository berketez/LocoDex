import os
import json
import hashlib
import sqlite3
from datetime import datetime
from typing import List, Dict, Any, Optional
import logging
from pathlib import Path

# Document processing libraries
import PyPDF2
import docx
import openpyxl
from pptx import Presentation

# Text processing
import re
from collections import Counter
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize

# Vector embeddings
import numpy as np
from sentence_transformers import SentenceTransformer

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class EnterpriseDocumentProcessor:
    """Kurumsal belge işleme ve indeksleme sistemi"""
    
    def __init__(self, db_path: str = "enterprise_documents.db"):
        self.db_path = db_path
        self.embedding_model = None
        self.setup_database()
        self.init_nlp()
        
        # Belge kategorileri ve anahtar kelimeleri
        self.categories = {
            'financial': {
                'name': 'Finansal Raporlar',
                'keywords': ['finansal', 'bütçe', 'gelir', 'gider', 'kâr', 'zarar', 'bilanço', 
                           'çeyrek', 'yıllık', 'revenue', 'profit', 'budget', 'finance'],
                'weight': 1.5
            },
            'hr': {
                'name': 'İnsan Kaynakları',
                'keywords': ['ik', 'personel', 'çalışan', 'oryantasyon', 'izin', 'maaş', 
                           'performans', 'işe alım', 'hr', 'employee', 'salary', 'hiring'],
                'weight': 1.3
            },
            'technical': {
                'name': 'Teknik Dokümanlar',
                'keywords': ['api', 'dokümantasyon', 'teknik', 'sistem', 'kurulum', 
                           'konfigürasyon', 'development', 'software', 'code', 'technical'],
                'weight': 1.2
            },
            'policy': {
                'name': 'Politika ve Prosedürler',
                'keywords': ['politika', 'prosedür', 'kural', 'yönetmelik', 'güvenlik', 
                           'uyumluluk', 'etik', 'policy', 'procedure', 'compliance'],
                'weight': 1.4
            },
            'marketing': {
                'name': 'Pazarlama ve Satış',
                'keywords': ['pazarlama', 'satış', 'kampanya', 'müşteri', 'analiz', 
                           'strateji', 'hedef', 'marketing', 'sales', 'customer'],
                'weight': 1.1
            },
            'training': {
                'name': 'Eğitim Materyalleri',
                'keywords': ['eğitim', 'kurs', 'öğrenim', 'sertifika', 'workshop', 
                           'seminer', 'gelişim', 'training', 'course', 'learning'],
                'weight': 1.0
            }
        }
        
        # Departman çıkarımı için anahtar kelimeler
        self.departments = {
            'İnsan Kaynakları': ['ik', 'hr', 'personel', 'çalışan', 'human resources'],
            'Finans': ['finans', 'muhasebe', 'bütçe', 'finance', 'accounting', 'budget'],
            'Pazarlama': ['pazarlama', 'satış', 'marketing', 'sales', 'müşteri', 'customer'],
            'Teknoloji': ['teknik', 'geliştirme', 'it', 'technical', 'development', 'software'],
            'Hukuk': ['hukuk', 'legal', 'uyumluluk', 'compliance', 'contract'],
            'Operasyon': ['operasyon', 'süreç', 'operation', 'process', 'logistics']
        }

    def setup_database(self):
        """Veritabanını kurulumu"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS documents (
                id TEXT PRIMARY KEY,
                filename TEXT NOT NULL,
                filepath TEXT NOT NULL,
                file_hash TEXT NOT NULL,
                content TEXT NOT NULL,
                summary TEXT,
                category TEXT,
                department TEXT,
                keywords TEXT,
                entities TEXT,
                upload_date TEXT,
                last_modified TEXT,
                file_size INTEGER,
                page_count INTEGER,
                language TEXT DEFAULT 'tr'
            )
        ''')
        
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS embeddings (
                document_id TEXT,
                chunk_id TEXT,
                chunk_text TEXT,
                embedding BLOB,
                PRIMARY KEY (document_id, chunk_id),
                FOREIGN KEY (document_id) REFERENCES documents (id)
            )
        ''')
        
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS search_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                query TEXT NOT NULL,
                results_count INTEGER,
                timestamp TEXT,
                response_time_ms INTEGER
            )
        ''')
        
        conn.commit()
        conn.close()

    def init_nlp(self):
        """NLP araçlarını başlat"""
        try:
            # NLTK veri setlerini indir
            nltk.download('punkt', quiet=True)
            nltk.download('stopwords', quiet=True)
            
            # Embedding modelini yükle
            self.embedding_model = SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')
            logger.info("NLP araçları başarıyla yüklendi")
        except Exception as e:
            logger.error(f"NLP araçları yüklenirken hata: {e}")

    def extract_text_from_file(self, filepath: str) -> Dict[str, Any]:
        """Dosyadan metin çıkarma"""
        try:
            file_ext = Path(filepath).suffix.lower()
            content = ""
            metadata = {"page_count": 0, "file_size": os.path.getsize(filepath)}
            
            if file_ext == '.pdf':
                content, metadata = self._extract_from_pdf(filepath)
            elif file_ext == '.docx':
                content, metadata = self._extract_from_docx(filepath)
            elif file_ext == '.xlsx':
                content, metadata = self._extract_from_xlsx(filepath)
            elif file_ext == '.pptx':
                content, metadata = self._extract_from_pptx(filepath)
            elif file_ext == '.txt':
                with open(filepath, 'r', encoding='utf-8') as file:
                    content = file.read()
                metadata["page_count"] = len(content.split('\n'))
            else:
                raise ValueError(f"Desteklenmeyen dosya türü: {file_ext}")
                
            return {"content": content, "metadata": metadata}
            
        except Exception as e:
            logger.error(f"Dosya işleme hatası {filepath}: {e}")
            raise

    def _extract_from_pdf(self, filepath: str) -> tuple:
        """PDF dosyasından metin çıkarma"""
        content = ""
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            page_count = len(pdf_reader.pages)
            
            for page in pdf_reader.pages:
                content += page.extract_text() + "\n"
                
        metadata = {"page_count": page_count, "file_size": os.path.getsize(filepath)}
        return content, metadata

    def _extract_from_docx(self, filepath: str) -> tuple:
        """DOCX dosyasından metin çıkarma"""
        doc = docx.Document(filepath)
        content = ""
        
        for paragraph in doc.paragraphs:
            content += paragraph.text + "\n"
            
        # Tablolardan da metin çıkar
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    content += cell.text + " "
                content += "\n"
                
        metadata = {"page_count": len(doc.paragraphs), "file_size": os.path.getsize(filepath)}
        return content, metadata

    def _extract_from_xlsx(self, filepath: str) -> tuple:
        """Excel dosyasından metin çıkarma"""
        workbook = openpyxl.load_workbook(filepath)
        content = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            content += f"Sheet: {sheet_name}\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    content += row_text + "\n"
                    
        metadata = {"page_count": len(workbook.sheetnames), "file_size": os.path.getsize(filepath)}
        return content, metadata

    def _extract_from_pptx(self, filepath: str) -> tuple:
        """PowerPoint dosyasından metin çıkarma"""
        presentation = Presentation(filepath)
        content = ""
        
        for i, slide in enumerate(presentation.slides):
            content += f"Slide {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
                    
        metadata = {"page_count": len(presentation.slides), "file_size": os.path.getsize(filepath)}
        return content, metadata

    def categorize_document(self, filename: str, content: str) -> str:
        """Belgeyi içeriğine göre kategorize et"""
        text = (filename + " " + content).lower()
        scores = {}
        
        for category, info in self.categories.items():
            score = 0
            for keyword in info['keywords']:
                score += text.count(keyword) * info['weight']
            scores[category] = score
            
        if max(scores.values()) > 0:
            return max(scores, key=scores.get)
        return 'general'

    def infer_department(self, filename: str, content: str) -> str:
        """Dosya içeriğinden departmanı çıkar"""
        text = (filename + " " + content).lower()
        
        for department, keywords in self.departments.items():
            for keyword in keywords:
                if keyword in text:
                    return department
        return 'Genel'

    def extract_keywords(self, content: str, top_k: int = 15) -> List[str]:
        """Metinden anahtar kelimeleri çıkar"""
        try:
            # Türkçe stopwords
            turkish_stopwords = set(['ve', 'ile', 'bu', 'şu', 'o', 'da', 'de', 'ki', 'mi', 'mu', 'mü', 
                                   'için', 'olan', 'olan', 'bir', 'çok', 'daha', 'en', 'gibi', 'kadar',
                                   'ancak', 'fakat', 'lakin', 'ama', 'sadece', 'yalnız', 'hem', 'ya', 
                                   'veya', 'yahut', 'yani', 'başka', 'diğer', 'bazı', 'her', 'hiç'])
            
            # İngilizce stopwords
            english_stopwords = set(stopwords.words('english'))
            all_stopwords = turkish_stopwords.union(english_stopwords)
            
            # Metni temizle ve tokenize et
            cleaned_text = re.sub(r'[^\w\s]', ' ', content.lower())
            words = word_tokenize(cleaned_text)
            
            # Filtreleme
            filtered_words = [
                word for word in words 
                if len(word) > 2 and word not in all_stopwords and word.isalpha()
            ]
            
            # Frekans analizi
            word_freq = Counter(filtered_words)
            return [word for word, freq in word_freq.most_common(top_k)]
            
        except Exception as e:
            logger.error(f"Anahtar kelime çıkarma hatası: {e}")
            return []

    def generate_summary(self, content: str, max_sentences: int = 3) -> str:
        """Belge özeti oluştur"""
        try:
            sentences = sent_tokenize(content)
            if len(sentences) <= max_sentences:
                return content[:500] + "..." if len(content) > 500 else content
            
            # En uzun ve anlamlı cümleleri seç
            sentence_scores = []
            for sentence in sentences:
                if len(sentence.split()) >= 5:  # En az 5 kelimeli cümleler
                    score = len(sentence.split())  # Basit scoring: kelime sayısı
                    sentence_scores.append((sentence, score))
            
            # En yüksek skorlu cümleleri al
            top_sentences = sorted(sentence_scores, key=lambda x: x[1], reverse=True)[:max_sentences]
            summary = ". ".join([sent[0] for sent in top_sentences])
            
            return summary if len(summary) <= 500 else summary[:500] + "..."
            
        except Exception as e:
            logger.error(f"Özet oluşturma hatası: {e}")
            return content[:300] + "..." if len(content) > 300 else content

    def create_embeddings(self, text: str, chunk_size: int = 512) -> List[np.ndarray]:
        """Metin için embedding vektörleri oluştur"""
        if not self.embedding_model:
            logger.warning("Embedding model yüklü değil")
            return []
            
        try:
            # Metni parçalara böl
            words = text.split()
            chunks = []
            
            for i in range(0, len(words), chunk_size):
                chunk = " ".join(words[i:i + chunk_size])
                chunks.append(chunk)
            
            # Her parça için embedding oluştur
            embeddings = self.embedding_model.encode(chunks)
            return embeddings
            
        except Exception as e:
            logger.error(f"Embedding oluşturma hatası: {e}")
            return []

    def process_document(self, filepath: str) -> Dict[str, Any]:
        """Belgeyi işle ve veritabanına kaydet"""
        try:
            # Dosya bilgilerini al
            filename = os.path.basename(filepath)
            file_hash = self._calculate_file_hash(filepath)
            
            # Dosya zaten işlenmiş mi kontrol et
            if self._is_document_processed(file_hash):
                logger.info(f"Belge zaten işlenmiş: {filename}")
                return {"status": "already_processed", "filename": filename}
            
            # Metni çıkar
            extraction_result = self.extract_text_from_file(filepath)
            content = extraction_result["content"]
            metadata = extraction_result["metadata"]
            
            # Analiz et
            category = self.categorize_document(filename, content)
            department = self.infer_department(filename, content)
            keywords = self.extract_keywords(content)
            summary = self.generate_summary(content)
            
            # Belge ID'si oluştur
            doc_id = hashlib.md5(f"{filename}_{file_hash}".encode()).hexdigest()
            
            # Veritabanına kaydet
            self._save_document_to_db(
                doc_id=doc_id,
                filename=filename,
                filepath=filepath,
                file_hash=file_hash,
                content=content,
                summary=summary,
                category=category,
                department=department,
                keywords=keywords,
                metadata=metadata
            )
            
            # Embeddings oluştur ve kaydet
            embeddings = self.create_embeddings(content)
            if embeddings:
                self._save_embeddings_to_db(doc_id, embeddings, content)
            
            logger.info(f"Belge başarıyla işlendi: {filename}")
            return {
                "status": "success",
                "doc_id": doc_id,
                "filename": filename,
                "category": category,
                "department": department,
                "summary": summary,
                "keywords": keywords[:5]  # İlk 5 anahtar kelime
            }
            
        except Exception as e:
            logger.error(f"Belge işleme hatası {filepath}: {e}")
            return {"status": "error", "error": str(e)}

    def _calculate_file_hash(self, filepath: str) -> str:
        """Dosya hash'i hesapla"""
        hash_md5 = hashlib.md5()
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def _is_document_processed(self, file_hash: str) -> bool:
        """Belge daha önce işlenmiş mi kontrol et"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        cursor.execute("SELECT id FROM documents WHERE file_hash = ?", (file_hash,))
        result = cursor.fetchone()
        conn.close()
        return result is not None

    def _save_document_to_db(self, doc_id: str, filename: str, filepath: str, 
                           file_hash: str, content: str, summary: str, 
                           category: str, department: str, keywords: List[str], 
                           metadata: Dict):
        """Belgeyi veritabanına kaydet"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute('''
            INSERT OR REPLACE INTO documents 
            (id, filename, filepath, file_hash, content, summary, category, 
             department, keywords, upload_date, last_modified, file_size, page_count)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ''', (
            doc_id, filename, filepath, file_hash, content, summary, category,
            department, json.dumps(keywords), datetime.now().isoformat(),
            datetime.now().isoformat(), metadata.get("file_size", 0),
            metadata.get("page_count", 0)
        ))
        
        conn.commit()
        conn.close()

    def _save_embeddings_to_db(self, doc_id: str, embeddings: List[np.ndarray], content: str):
        """Embeddings'i veritabanına kaydet"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        # Önce mevcut embeddings'i sil
        cursor.execute("DELETE FROM embeddings WHERE document_id = ?", (doc_id,))
        
        # Yeni embeddings'i kaydet
        words = content.split()
        chunk_size = 512
        
        for i, embedding in enumerate(embeddings):
            chunk_start = i * chunk_size
            chunk_end = min((i + 1) * chunk_size, len(words))
            chunk_text = " ".join(words[chunk_start:chunk_end])
            
            cursor.execute('''
                INSERT INTO embeddings (document_id, chunk_id, chunk_text, embedding)
                VALUES (?, ?, ?, ?)
            ''', (doc_id, f"chunk_{i}", chunk_text, embedding.tobytes()))
        
        conn.commit()
        conn.close()

    def search_documents(self, query: str, top_k: int = 5, category_filter: str = None) -> List[Dict]:
        """Semantik arama yap"""
        start_time = datetime.now()
        
        try:
            # Query embedding'i oluştur
            if self.embedding_model:
                query_embedding = self.embedding_model.encode([query])[0]
                semantic_results = self._semantic_search(query_embedding, top_k * 2)
            else:
                semantic_results = []
            
            # Keyword-based arama
            keyword_results = self._keyword_search(query, top_k * 2, category_filter)
            
            # Sonuçları birleştir ve skorla
            all_results = {}
            
            # Semantic sonuçları ekle
            for result in semantic_results:
                doc_id = result["doc_id"]
                if doc_id not in all_results:
                    all_results[doc_id] = result
                    all_results[doc_id]["semantic_score"] = result["score"]
                else:
                    all_results[doc_id]["semantic_score"] = max(
                        all_results[doc_id].get("semantic_score", 0), result["score"]
                    )
            
            # Keyword sonuçları ekle
            for result in keyword_results:
                doc_id = result["doc_id"]
                if doc_id not in all_results:
                    all_results[doc_id] = result
                    all_results[doc_id]["keyword_score"] = result["score"]
                else:
                    all_results[doc_id]["keyword_score"] = result["score"]
            
            # Final skoru hesapla
            final_results = []
            for doc_id, result in all_results.items():
                semantic_score = result.get("semantic_score", 0)
                keyword_score = result.get("keyword_score", 0)
                final_score = (semantic_score * 0.6) + (keyword_score * 0.4)
                
                result["final_score"] = final_score
                final_results.append(result)
            
            # Skora göre sırala ve top-k al
            final_results.sort(key=lambda x: x["final_score"], reverse=True)
            top_results = final_results[:top_k]
            
            # Arama geçmişini kaydet
            response_time = (datetime.now() - start_time).total_seconds() * 1000
            self._save_search_history(query, len(top_results), response_time)
            
            return top_results
            
        except Exception as e:
            logger.error(f"Arama hatası: {e}")
            return []

    def _semantic_search(self, query_embedding: np.ndarray, top_k: int) -> List[Dict]:
        """Semantik arama yap"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT e.document_id, e.chunk_text, e.embedding, d.filename, d.category, d.department, d.summary
            FROM embeddings e
            JOIN documents d ON e.document_id = d.id
        ''')
        
        results = []
        for row in cursor.fetchall():
            doc_id, chunk_text, embedding_bytes, filename, category, department, summary = row
            
            # Embedding'i numpy array'e çevir
            chunk_embedding = np.frombuffer(embedding_bytes, dtype=np.float32)
            
            # Cosine similarity hesapla
            similarity = np.dot(query_embedding, chunk_embedding) / (
                np.linalg.norm(query_embedding) * np.linalg.norm(chunk_embedding)
            )
            
            if similarity > 0.3:  # Minimum similarity threshold
                results.append({
                    "doc_id": doc_id,
                    "filename": filename,
                    "category": category,
                    "department": department,
                    "summary": summary,
                    "matched_text": chunk_text[:200] + "...",
                    "score": float(similarity)
                })
        
        conn.close()
        return sorted(results, key=lambda x: x["score"], reverse=True)[:top_k]

    def _keyword_search(self, query: str, top_k: int, category_filter: str = None) -> List[Dict]:
        """Anahtar kelime tabanlı arama"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        query_lower = query.lower()
        query_words = query_lower.split()
        
        # SQL sorgusu hazırla
        sql = '''
            SELECT id, filename, content, category, department, summary, keywords
            FROM documents
        '''
        params = []
        
        if category_filter and category_filter != 'all':
            sql += ' WHERE category = ?'
            params.append(category_filter)
        
        cursor.execute(sql, params)
        
        results = []
        for row in cursor.fetchall():
            doc_id, filename, content, category, department, summary, keywords_json = row
            
            content_lower = content.lower()
            filename_lower = filename.lower()
            
            score = 0
            
            # Dosya adında arama (yüksek puan)
            for word in query_words:
                score += filename_lower.count(word) * 10
            
            # İçerikte arama
            for word in query_words:
                score += content_lower.count(word) * 2
            
            # Anahtar kelimelerde arama
            try:
                keywords = json.loads(keywords_json)
                for keyword in keywords:
                    for word in query_words:
                        if word in keyword.lower():
                            score += 5
            except:
                pass
            
            if score > 0:
                # Matched text çıkar
                matched_text = ""
                for sentence in content.split('.'):
                    if any(word in sentence.lower() for word in query_words):
                        matched_text = sentence[:200] + "..."
                        break
                
                results.append({
                    "doc_id": doc_id,
                    "filename": filename,
                    "category": category,
                    "department": department,
                    "summary": summary,
                    "matched_text": matched_text,
                    "score": score
                })
        
        conn.close()
        return sorted(results, key=lambda x: x["score"], reverse=True)[:top_k]

    def _save_search_history(self, query: str, results_count: int, response_time_ms: float):
        """Arama geçmişini kaydet"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute('''
            INSERT INTO search_history (query, results_count, timestamp, response_time_ms)
            VALUES (?, ?, ?, ?)
        ''', (query, results_count, datetime.now().isoformat(), int(response_time_ms)))
        
        conn.commit()
        conn.close()

    def get_document_stats(self) -> Dict[str, Any]:
        """Belge istatistiklerini al"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        # Toplam belge sayısı
        cursor.execute("SELECT COUNT(*) FROM documents")
        total_docs = cursor.fetchone()[0]
        
        # Kategoriye göre dağılım
        cursor.execute("SELECT category, COUNT(*) FROM documents GROUP BY category")
        category_stats = dict(cursor.fetchall())
        
        # Departmana göre dağılım
        cursor.execute("SELECT department, COUNT(*) FROM documents GROUP BY department")
        department_stats = dict(cursor.fetchall())
        
        # Son 7 günün arama istatistikleri
        cursor.execute('''
            SELECT COUNT(*), AVG(response_time_ms)
            FROM search_history 
            WHERE timestamp >= datetime('now', '-7 days')
        ''')
        search_stats = cursor.fetchone()
        
        conn.close()
        
        return {
            "total_documents": total_docs,
            "categories": category_stats,
            "departments": department_stats,
            "recent_searches": search_stats[0] or 0,
            "avg_response_time_ms": search_stats[1] or 0
        }

    def get_document_by_id(self, doc_id: str) -> Optional[Dict[str, Any]]:
        """ID'ye göre belge getir"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute("SELECT * FROM documents WHERE id = ?", (doc_id,))
        row = cursor.fetchone()
        
        if row:
            columns = [description[0] for description in cursor.description]
            document = dict(zip(columns, row))
            
            # Keywords'ü parse et
            try:
                document['keywords'] = json.loads(document['keywords'])
            except:
                document['keywords'] = []
                
            conn.close()
            return document
        
        conn.close()
        return None

# CLI kullanımı için yardımcı fonksiyonlar
def process_directory(directory_path: str, processor: EnterpriseDocumentProcessor):
    """Bir dizindeki tüm belgeleri işle"""
    supported_extensions = {'.pdf', '.docx', '.xlsx', '.pptx', '.txt'}
    processed_count = 0
    
    for root, dirs, files in os.walk(directory_path):
        for file in files:
            if Path(file).suffix.lower() in supported_extensions:
                filepath = os.path.join(root, file)
                try:
                    result = processor.process_document(filepath)
                    if result["status"] == "success":
                        processed_count += 1
                        print(f"✅ İşlendi: {file} - {result['category']} - {result['department']}")
                    elif result["status"] == "already_processed":
                        print(f"⏭️  Zaten işlenmiş: {file}")
                    else:
                        print(f"❌ Hata: {file} - {result.get('error', 'Bilinmeyen hata')}")
                except Exception as e:
                    print(f"❌ Kritik hata: {file} - {e}")
    
    print(f"\n📊 Toplam {processed_count} belge işlendi")

if __name__ == "__main__":
    # Test için basit kullanım
    processor = EnterpriseDocumentProcessor()
    
    # Örnek arama
    results = processor.search_documents("finansal rapor", top_k=5)
    print("Arama sonuçları:")
    for result in results:
        print(f"- {result['filename']} (Skor: {result['final_score']:.2f})")